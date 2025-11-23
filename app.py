import os
import json
import pickle
import requests
import torch
import faiss
import numpy as np
from tqdm.auto import tqdm
from spacy.lang.en import English
from sentence_transformers import SentenceTransformer
import fitz
import docx
from pptx import Presentation
import streamlit as st
import subprocess
import time


# =======================
#  Ollama LLM Interface
# =======================
def start_ollama_server():
    try:
        response = requests.get("http://localhost:11434/api/tags", timeout=2)
        if response.status_code == 200:
            print("✅ Ollama is already running.")
            return
    except requests.exceptions.RequestException:
        print("⚠️ Ollama not running. Starting Ollama server...")

    subprocess.Popen(
        ["ollama", "serve"],
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        creationflags=subprocess.CREATE_NEW_CONSOLE if os.name == "nt" else 0
    )

    print("⏳ Waiting for Ollama to start...")
    time.sleep(5)

    try:
        response = requests.get("http://localhost:11434/api/tags", timeout=5)
        if response.status_code == 200:
            print("✅ Ollama started successfully.")
    except Exception:
        print("❌ Failed to connect to Ollama after starting.")


def query_ollama(prompt: str, model: str = "gemma:2b") -> str:
    try:
        response = requests.post(
            "http://localhost:11434/api/generate",
            json={"model": model, "prompt": prompt},
            stream=True,
        )
        response.raise_for_status()
    except requests.exceptions.ConnectionError:
        return "❌ Ollama server not running. Start it first with: ollama serve"
    except Exception as e:
        return f"❌ Error querying Ollama: {e}"

    output = ""
    for line in response.iter_lines():
        if line:
            data = json.loads(line)
            output += data.get("response", "")
    return output.strip()


# =======================
#  Document Processing
# =======================
def process_document(doc_name, output_dir="embeddings_store"):
    base_name = os.path.splitext(os.path.basename(doc_name))[0]
    pkl_path = os.path.join(output_dir, f"{base_name}.pkl")

    if os.path.exists(pkl_path):
        with open(pkl_path, "rb") as f:
            with st.spinner(f"📄 '{base_name}' already processed — loading from disk..."):
                return pickle.load(f)

    if doc_name.endswith("pdf"):
        doc = open_and_read_pdf(doc_name)
    elif doc_name.endswith("pptx"):
        doc = open_and_read_pptx(doc_name)
    elif doc_name.endswith("docx"):
        doc = open_and_read_docx(doc_name)
    else:
        doc = open_and_read_txt(doc_name)

    with st.spinner(f"Generating Sentences & Chunks for '{base_name}'..."):
        doc = generate_sentences_and_chunks(doc)
    with st.spinner(f"Generating Embeddings for '{base_name}'..."):
        doc = generate_embeddings(doc)
    with st.spinner(f"Flattening & Saving Embeddings for '{base_name}'..."):
        records = flatten_embeddings(doc, doc_name)
        save_embeddings_per_file(records, doc_name, output_dir)

    return records


# =======================
#  RAG Core Functions
# =======================
def generate_answer(query: str, retrieved_chunks: list[dict], model: str = "gemma:2b"):
    context = " | ".join([chunk["text"] for chunk in retrieved_chunks])
    prompt = (
        f"Act as an expert assistant. Given the query: '{query}' "
        f"and the retrieved chunks: '{context}', answer strictly using only this context in short and summarized answer."
        f"If the retrieved chunks do not answer the question, Say the passages are irrelevant"
    )
    with st.spinner(f"Generating Answer ..."):
        llm_answer = query_ollama(prompt, model)
        sources = [f"{c.get('source_file', '?')} (p.{c.get('page_number', '?')})" for c in retrieved_chunks]
        sources_text = " | ".join(sorted(set(sources))) if sources else "No sources found."
        return f"{llm_answer}\n\n Sources: {sources_text}"
   


def retrieve(query: str, records, faiss_index, top_k=3, model_name="BAAI/bge-large-en-v1.5"):
    with st.spinner(f"Retrieving Relevant Chunks ...."):
        if faiss_index is None:
            return []
        query_vector = embed_query(query, model_name=model_name)
        D, I = faiss_index.search(np.array([query_vector]), top_k)
        results = [records[int(idx)] for idx in I[0] if idx >= 0 and idx < len(records)]
        return results


# =======================
#  Utilities
# =======================
def text_formatter(text: str) -> str:
    return text.replace("\n", " ").strip()


def open_and_read_pdf(path):
    doc = fitz.open(path)
    return [{"page_number": i + 1, "text": text_formatter(p.get_text())} for i, p in enumerate(doc)]


def open_and_read_docx(path):
    document = docx.Document(path)
    return [{"page_number": i + 1, "text": text_formatter(p.text)} for i, p in enumerate(document.paragraphs) if p.text.strip()]


def open_and_read_pptx(path):
    pres = Presentation(path)
    slides = []
    for i, slide in enumerate(pres.slides):
        text = " ".join(shape.text for shape in slide.shapes if hasattr(shape, "text"))
        if text.strip():
            slides.append({"page_number": i + 1, "text": text_formatter(text)})
    return slides


def open_and_read_txt(path):
    with open(path, "r", encoding="utf-8") as f:
        text = f.read()
    return [{"page_number": 1, "text": text_formatter(text)}]


def generate_sentences_and_chunks(doc, chunk_size=550, overlap=100):
    nlp = English()
    nlp.add_pipe("sentencizer")
    for item in doc:
        words = item["text"].split()
        chunks = [" ".join(words[i:i + chunk_size]) for i in range(0, len(words), chunk_size - overlap)]
        item["chunks"] = chunks
    return doc


def generate_embeddings(doc, model_name="BAAI/bge-large-en-v1.5", device=None, batch_size=32):
    device = device or ("cuda" if torch.cuda.is_available() else "cpu")
    model = SentenceTransformer(f"./{model_name.split('/')[-1]}", device=device)
    for item in doc:
        text_units = item.get("chunks", [])
        if not text_units:
            item["embeddings"] = []
            continue
        embeddings = model.encode(text_units, batch_size=batch_size, convert_to_numpy=True, normalize_embeddings=True)
        item["embeddings"] = embeddings
        item["embedded_texts"] = text_units
    return doc

def save_embeddings_per_file(records, file_name, output_dir="embeddings_store"):
    os.makedirs(output_dir, exist_ok=True)
    base_name = os.path.splitext(os.path.basename(file_name))[0]
    path = os.path.join(output_dir, f"{base_name}.pkl")
    with open(path, "wb") as f:
        pickle.dump(records, f)
    return path

def embed_query(query, model_name="BAAI/bge-large-en-v1.5", device=None):
    device = device or ("cuda" if torch.cuda.is_available() else "cpu")
    model = SentenceTransformer(f"./{model_name.split('/')[-1]}", device=device)
    return model.encode([query], convert_to_numpy=True, normalize_embeddings=True)[0]


def flatten_embeddings(doc, file_path):
    records = []
    base_name = os.path.splitext(os.path.basename(file_path))[0]
    for item in doc:
        for text, vector in zip(item.get("embedded_texts", []), item.get("embeddings", [])):
            records.append({
                "text": text,
                "embedding": vector,
                "page_number": item.get("page_number", "?"),
                "source_file": base_name
            })
    return records




def build_faiss_index(records):
    if not records:
        return None
    with st.spinner(f"Building FAISS Index... "):
        dim = len(records[0]["embedding"])
        index = faiss.IndexFlatIP(dim)
        vectors = np.vstack([r["embedding"] for r in records])
        index.add(vectors)
        return index


# =======================
#  Streamlit Frontend
# =======================
st.title("🔍 Local RAG Assistant")
start_ollama_server()

if "chat_history" not in st.session_state:
    st.session_state["chat_history"] = []

st.subheader("Conversation History")
if st.session_state["chat_history"]:
    for msg in st.session_state["chat_history"]:
        role = msg["role"]
        st.markdown(f"**{'You' if role == 'user' else 'Assistant'}:** {msg['message']}")
else:
    st.info("No messages yet — upload a document and ask a question to start.")

st.markdown("---")
st.subheader("📂 Upload Documents")
uploaded_files = st.file_uploader(
    "Upload one or more documents",
    type=["pdf", "docx", "pptx", "txt"],
    accept_multiple_files=True
)

os.makedirs("temp_docs", exist_ok=True)
os.makedirs("embeddings_store", exist_ok=True)

combined_records = []

if uploaded_files:
    for uploaded_file in uploaded_files:
        temp_path = os.path.join("temp_docs", uploaded_file.name)
        with open(temp_path, "wb") as f:
            f.write(uploaded_file.getbuffer())

        try:
            recs = process_document(temp_path)
            combined_records.extend(recs)
        except Exception as e:
            st.error(f"Error processing {uploaded_file.name}: {e}")

    st.success(f"✅ Processed {len(uploaded_files)} file(s). Total chunks: {len(combined_records)}")
else:
    st.info("Upload at least one document to enable retrieval.")

st.markdown("---")
query = st.text_input("Ask your question:")

if st.button("Generate Answer") and query:
    if not combined_records:
        st.warning("Please upload and process at least one document first!")
    else:
        try:
            index = build_faiss_index(combined_records)
            retrieved = retrieve(query, combined_records, index, top_k=3)
            answer = generate_answer(query, retrieved)

            st.session_state["chat_history"].append({"role": "user", "message": query})
            st.session_state["chat_history"].append({"role": "assistant", "message": answer})

            st.subheader("🤖 Answer")
            st.write(answer)

            st.subheader("📚 Retrieved Chunks")
            for r in retrieved:
                st.markdown(f"- **{r['source_file']} (p.{r['page_number']})** — {r['text'][:200]}...")

        except Exception as e:
            st.error(f"Error generating answer: {e}")

if st.button("🧹 Clear conversation history"):
    st.session_state["chat_history"] = []
    st.success("Conversation history cleared.")
